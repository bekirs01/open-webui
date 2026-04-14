"""
MWS Auto: server-side presentation generation (structured plan → AI images → PPTX file).

HTML is not the source of truth; the deck is built with python-pptx so layout is stable
and each slide image is generated from an explicit English image_prompt tied to that slide.
"""

from __future__ import annotations

import io
import json
import logging
import os
import re
from typing import Any

from fastapi import Request, UploadFile

log = logging.getLogger(__name__)

# Env: MWS_PRESENTATION_PIPELINE=false disables this path (LLM answers normally).
# MWS_PRESENTATION_AUTO_ONLY=false allows manual model selection + presentation intent.


def _pipeline_enabled() -> bool:
    return os.environ.get('MWS_PRESENTATION_PIPELINE', 'true').lower() == 'true'


def _auto_only() -> bool:
    return os.environ.get('MWS_PRESENTATION_AUTO_ONLY', 'false').lower() == 'true'


def _desired_slide_count(message_text: str) -> int:
    t = message_text or ''
    short = bool(re.search(r'\b(?:коротк|кратк|небольш|мало\s+слайд|few\s+slide|short)\b', t, re.I))
    for pat in (
        r'(\d+)\s*(?:слайд|слайдов|слайда|slide|slides|slayt)',
        r'(?:слайд|slide|slides)\s*[:\-]?\s*(\d+)',
    ):
        m = re.search(pat, t, re.I)
        if m:
            n = max(2, min(int(m.group(1)), 15))
            return n
    if short:
        return 5
    return 7


def _guess_lang_reply(message_text: str) -> str:
    t = message_text or ''
    if re.search(r'[а-яА-ЯёЁ]', t):
        return 'ru'
    if re.search(r'[ğüşöçıİĞÜŞÖÇ]', t, re.I):
        return 'tr'
    return 'en'


def _done_message(lang: str) -> str:
    if lang == 'ru':
        return 'Презентация PowerPoint (PPTX) готова — файл прикреплён, можно открыть и скачать.'
    if lang == 'tr':
        return 'PowerPoint sunumu (PPTX) hazır; dosya eklendi, indirebilirsiniz.'
    return 'Your PowerPoint presentation (PPTX) is attached — open or download the file below.'


def _build_error_message(lang: str) -> str:
    if lang == 'ru':
        return (
            'Не удалось собрать презентацию (ошибка при сборке или загрузке файла). '
            'Попробуйте снова или укажите меньше слайдов.'
        )
    if lang == 'tr':
        return (
            'Sunum oluşturulamadı (derleme veya dosya yükleme hatası). '
            'Daha az slayt ile tekrar deneyin.'
        )
    return (
        'Could not build the presentation (build or upload failed). '
        'Try again or ask for fewer slides.'
    )


def _extract_json_object(text: str) -> dict[str, Any] | None:
    if not text:
        return None
    s = text.strip()
    # fenced ```json ... ```
    try:
        m = re.search(r'```(?:json)?\s*(\{[\s\S]*?\})\s*```', s, re.I)
        if m:
            data = json.loads(m.group(1))
            if isinstance(data, dict):
                return data
    except Exception:
        pass
    try:
        bracket_start = s.find('{')
        bracket_end = s.rfind('}') + 1
        if bracket_start == -1 or bracket_end <= bracket_start:
            return None
        blob = s[bracket_start:bracket_end]
        data = json.loads(blob)
        return data if isinstance(data, dict) else None
    except Exception:
        return None


def _planner_failed_form_data(form_data: dict[str, Any], user_text: str) -> dict[str, Any]:
    lang = _guess_lang_reply(user_text)
    if lang == 'ru':
        msg = (
            'Не удалось автоматически собрать презентацию: модель не вернула корректный JSON-план. '
            'Повторите запрос короче («создай презентацию про …, 5 слайдов») или попробуйте снова через минуту.'
        )
    elif lang == 'tr':
        msg = (
            'Sunum oluşturulamadı: planlayıcı geçerli JSON döndürmedi. '
            'Daha kısa bir istekle tekrar deneyin.'
        )
    else:
        msg = (
            'Could not build the presentation: the planner did not return valid JSON. '
            'Retry with a shorter request (e.g. topic + slide count).'
        )
    form_data['_mws_export_completion'] = True
    form_data['_mws_export_assistant_content'] = msg
    form_data['features'] = {}
    return form_data


_PLANNER_RETRY_SYSTEM = """You MUST output a single valid JSON object only. No markdown, no commentary.

Required shape:
{"title":"string","slides":[{"title":"string","bullets":["string"],"image_prompt":"English image description"}]}

The "slides" array length must equal N from the user message. image_prompt must be English and on-topic for that slide."""


def _read_uploaded_image_bytes(url: str) -> bytes | None:
    """Load bytes for a file served via /files/{id}/content (same user upload)."""
    try:
        m = re.search(r'/files/([^/]+)/', url or '')
        if not m:
            return None
        from pathlib import Path

        from open_webui.models.files import Files
        from open_webui.storage.provider import Storage

        fid = m.group(1)
        fm = Files.get_file_by_id(fid)
        if not fm or not fm.path:
            return None
        p = Storage.get_file(fm.path)
        path = Path(p)
        if path.is_file():
            return path.read_bytes()
    except Exception as e:
        log.debug('[MWS] presentation read image bytes: %s', e)
    return None


def _slug_filename(title: str) -> str:
    raw = (title or 'presentation').strip()
    s = re.sub(r'[^\w\s\-]', '', raw, flags=re.UNICODE)
    s = re.sub(r'[\s\-]+', '-', s).strip('-').lower()
    return (s[:56] or 'presentation') + '.pptx'


def _build_pptx_bytes(
    *,
    slides: list[dict[str, Any]],
    images: list[bytes | None],
) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    blank = prs.slide_layouts[6]

    for i, slide in enumerate(slides):
        title = (slide.get('title') or f'Slide {i + 1}').strip()[:200]
        bullets = slide.get('bullets') if isinstance(slide.get('bullets'), list) else []
        bullets = [str(b).strip()[:500] for b in bullets if str(b).strip()][:6]

        s = prs.slides.add_slide(blank)

        # Title
        tbox = s.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.85))
        tf = tbox.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(26)

        img_bytes = images[i] if i < len(images) else None
        if img_bytes:
            try:
                s.shapes.add_picture(io.BytesIO(img_bytes), Inches(0.5), Inches(1.25), width=Inches(6.2))
            except Exception as e:
                log.debug('[MWS] presentation slide image skip: %s', e)

        # Bullets (right column, or full width if no image)
        left_text = Inches(6.95) if img_bytes else Inches(0.55)
        tw = Inches(6.0) if img_bytes else Inches(12.2)
        bbox = s.shapes.add_textbox(left_text, Inches(1.25), tw, Inches(5.6))
        bf = bbox.text_frame
        bf.word_wrap = True
        if bullets:
            bf.text = bullets[0]
            bf.paragraphs[0].font.size = Pt(17)
            for b in bullets[1:]:
                p = bf.add_paragraph()
                p.text = b
                p.level = 0
                p.font.size = Pt(17)
        else:
            bf.text = ' '

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


_PLANNER_SYSTEM = """You are an expert presentation planner. Output ONE JSON object only (no markdown fences).

Schema:
{
  "title": "short deck title in the user's language",
  "slides": [
    {
      "title": "slide title in the user's language",
      "bullets": ["...", "..."],
      "image_prompt": "Single detailed English prompt for an illustration/photo-style image that matches ONLY this slide. No text to render in the image. Consistent with the deck topic."
    }
  ]
}

Rules:
- slides length must be exactly N (given in the user message).
- Each image_prompt must describe visuals for that slide only (no random unrelated topics).
- Bullets: 2–5 items per slide, concise, in the SAME language as the user message (natural sentences).
- Do NOT output broken text, mixed scripts, stray punctuation, braces, or English gibberish inside bullets.
- image_prompt: English only, concrete, safe for work, no words/logos to draw as text."""


async def try_mws_presentation_pipeline(
    request: Request,
    form_data: dict[str, Any],
    user: Any,
    metadata: dict[str, Any],
) -> dict[str, Any]:
    """
    If conditions match, build PPTX with AI images and attach; set _mws_export_completion like export_pipeline.
    """
    try:
        from open_webui.utils.mws_gpt.active import is_mws_gpt_active
        from open_webui.utils.mws_gpt.presentation_intent import resolve_presentation_intent
        from open_webui.utils.mws_gpt.registry import extract_last_user_text
        from open_webui.utils.mws_gpt.router import LEGACY_AUTO_IDS

        if not _pipeline_enabled() or not is_mws_gpt_active(request.app.state.config):
            return form_data
        if form_data.get('_mws_export_completion'):
            return form_data
        if metadata.get('params', {}).get('function_calling') == 'native':
            return form_data

        user_text = extract_last_user_text(form_data.get('messages') or [])
        if not resolve_presentation_intent(user_text):
            return form_data

        mws_meta = metadata.get('mws_routing') or {}
        if _auto_only():
            orig = (mws_meta.get('original_requested_id') or '').strip()
            # Пустой orig: клиент/UI уже подставил модель без meta — не блокируем (legacy Auto).
            if orig and orig not in LEGACY_AUTO_IDS:
                return form_data

        chat_id = metadata.get('chat_id')
        message_id = metadata.get('message_id')
        if not chat_id or str(chat_id).startswith('local:') or not message_id:
            return form_data

        cfg = request.app.state.config
        if not getattr(cfg, 'ENABLE_IMAGE_GENERATION', False):
            form_data['_mws_export_completion'] = True
            form_data['_mws_export_assistant_content'] = (
                'Включите генерацию изображений в настройках (ENABLE_IMAGE_GENERATION), чтобы собрать слайды с картинками.'
                if _guess_lang_reply(user_text) == 'ru'
                else 'Enable image generation in settings to build slides with AI images.'
            )
            form_data['features'] = {}
            return form_data

        models = request.app.state.OPENAI_MODELS or {}
        from open_webui.utils.mws_gpt.team_registry import (
            filter_team_available,
            first_available,
            AUTO_IMAGE_ORDER,
        )

        available = filter_team_available(set(models.keys()))
        image_model_id = first_available(AUTO_IMAGE_ORDER, available)
        if not image_model_id:
            log.warning('[MWS] presentation: no image model in team set')
            form_data['_mws_export_completion'] = True
            form_data['_mws_export_assistant_content'] = (
                'Нет доступной модели генерации изображений (нужна из команды, напр. qwen-image).'
                if _guess_lang_reply(user_text) == 'ru'
                else 'No image generation model is available for this workspace.'
            )
            form_data['features'] = {}
            return form_data

        text_model_id = (form_data.get('model') or '').strip()
        if not text_model_id:
            return form_data

        n_slides = _desired_slide_count(user_text)
        planner_user = (
            f'User request:\n{user_text}\n\n'
            f'Produce exactly {n_slides} objects in the "slides" array (see system schema).'
        )

        mini_meta = {
            k: metadata[k]
            for k in ('chat_id', 'session_id', 'message_id', 'parent_message_id')
            if k in metadata
        }

        from open_webui.routers.openai import generate_chat_completion as openai_generate_chat_completion

        inner: dict[str, Any] = {
            'model': text_model_id,
            'messages': [
                {'role': 'system', 'content': _PLANNER_SYSTEM},
                {'role': 'user', 'content': planner_user},
            ],
            'stream': False,
            'metadata': mini_meta,
        }
        resp = await openai_generate_chat_completion(request, inner, user)
        raw = ''
        if isinstance(resp, dict):
            ch = (resp.get('choices') or [{}])[0]
            msg = ch.get('message') or {}
            raw = (msg.get('content') or '').strip()

        plan = _extract_json_object(raw)
        if not plan or not isinstance(plan.get('slides'), list):
            log.warning('[MWS] presentation: planner retry (strict JSON)')
            inner_retry: dict[str, Any] = {
                'model': text_model_id,
                'messages': [
                    {'role': 'system', 'content': _PLANNER_RETRY_SYSTEM},
                    {
                        'role': 'user',
                        'content': f'{planner_user}\n\nN must be exactly {n_slides}. Reply with JSON only, one object.',
                    },
                ],
                'stream': False,
                'metadata': mini_meta,
            }
            resp2 = await openai_generate_chat_completion(request, inner_retry, user)
            raw2 = ''
            if isinstance(resp2, dict):
                ch2 = (resp2.get('choices') or [{}])[0]
                msg2 = ch2.get('message') or {}
                raw2 = (msg2.get('content') or '').strip()
            plan = _extract_json_object(raw2)

        if not plan or not isinstance(plan.get('slides'), list):
            log.warning('[MWS] presentation: bad planner JSON after retry')
            return _planner_failed_form_data(form_data, user_text)

        slides_raw = [s for s in plan['slides'] if isinstance(s, dict)]
        if len(slides_raw) < 1:
            return _planner_failed_form_data(form_data, user_text)
        slides_raw = slides_raw[:n_slides]

        from open_webui.socket.main import get_event_emitter

        emitter = get_event_emitter(metadata)
        if emitter:
            await emitter({'type': 'status', 'data': {'description': 'Building presentation', 'done': False}})

        prev_img_model = getattr(cfg, 'IMAGE_GENERATION_MODEL', None)
        try:
            from open_webui.routers.images import (
                CreateImageForm,
                image_generations,
                set_image_model,
            )
            from open_webui.utils.mws_gpt.image_prompt import get_default_image_negative_prompt

            set_image_model(request, image_model_id)

            slide_images: list[bytes | None] = []
            deck_topic = (plan.get('title') or '').strip()[:200]

            for idx, sd in enumerate(slides_raw):
                ip = (sd.get('image_prompt') or '').strip()
                if not ip:
                    ip = (
                        f'Professional presentation illustration for slide titled "{sd.get("title", "")}" '
                        f'about {deck_topic}. Clean, modern, coherent with the deck topic.'
                    )
                ip = f'{ip}\n\nTopic context (do not render as text): {deck_topic}.'

                # Slide images are embedded into PPTX only. Do NOT pass chat_id/message_id here:
                # upload_image() attaches to the assistant message whenever both are set, which made
                # the UI show raw generated images instead of a single .pptx.
                img_meta: dict[str, Any] = {'mws_skip_message_image_attach': True}
                _sz = (os.environ.get('MWS_PRESENTATION_IMAGE_SIZE') or '1024x1024').strip()
                np = get_default_image_negative_prompt()
                cf: dict[str, Any] = {'prompt': ip, 'n': 1}
                if _sz and 'x' in _sz:
                    cf['size'] = _sz
                if np and (os.environ.get('MWS_IMAGE_NEGATIVE_PROMPT', 'true') or 'true').lower() != 'false':
                    cf['negative_prompt'] = np
                img_form = CreateImageForm(**cf)

                try:
                    imgs = await image_generations(
                        request=request,
                        form_data=img_form,
                        metadata=img_meta,
                        user=user,
                    )
                    url = (imgs[0].get('url') or '') if imgs else ''
                    b = _read_uploaded_image_bytes(url) if url else None
                    slide_images.append(b)
                except Exception as e:
                    log.warning('[MWS] presentation slide %s image failed: %s', idx, e)
                    slide_images.append(None)

            pptx_bytes = _build_pptx_bytes(slides=slides_raw, images=slide_images)
            if not pptx_bytes or len(pptx_bytes) < 400:
                raise ValueError('empty_pptx')

            fname = _slug_filename(plan.get('title') or 'presentation')
            from open_webui.models.chats import Chats
            from open_webui.routers.files import upload_file_handler

            uf = UploadFile(
                file=io.BytesIO(pptx_bytes),
                filename=fname,
                headers={'content-type': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'},
            )
            file_item = upload_file_handler(
                request,
                file=uf,
                metadata={'chat_id': chat_id, 'message_id': message_id},
                process=False,
                user=user,
            )
            if not file_item or not getattr(file_item, 'id', None):
                raise ValueError('upload_failed')

            try:
                Chats.insert_chat_files(
                    chat_id=chat_id,
                    message_id=message_id,
                    file_ids=[file_item.id],
                    user_id=user.id,
                )
            except Exception as e:
                log.debug('[MWS] presentation insert_chat_files: %s', e)

            lang = _guess_lang_reply(user_text)
            from open_webui.utils.mws_gpt.export_pipeline import _export_file_attachment_record

            ct = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            entry = _export_file_attachment_record(file_item, pptx_bytes, fname, ct)
            form_data['_mws_export_result_files'] = [entry]

            if emitter:
                await emitter({'type': 'status', 'data': {'description': 'Done', 'done': True}})
                await emitter({'type': 'files', 'data': {'files': [entry]}})

            form_data['_mws_export_completion'] = True
            form_data['_mws_export_assistant_content'] = _done_message(lang)
            form_data['features'] = {}
            return form_data

        except Exception as e:
            log.warning('[MWS] presentation build/upload failed: %s', e)
            lang = _guess_lang_reply(user_text)
            if emitter:
                try:
                    await emitter(
                        {
                            'type': 'status',
                            'data': {'description': 'Presentation failed', 'done': True},
                        }
                    )
                except Exception:
                    pass
            form_data['_mws_export_completion'] = True
            form_data['_mws_export_assistant_content'] = _build_error_message(lang)
            form_data['features'] = {}
            return form_data
        finally:
            try:
                cfg.IMAGE_GENERATION_MODEL = prev_img_model
            except Exception:
                pass

    except Exception as e:
        log.warning('[MWS] presentation pipeline: %s', e)
        return form_data

    return form_data
